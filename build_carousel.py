"""
Build a 32:9 carousel PowerPoint presentation from slide images.

This script creates a carousel-style presentation where multiple slide images
are arranged on each page in a grid layout.

Usage:
    python build_carousel.py images_folder/ output_carousel.pptx --slides-per-page 4
"""

import argparse
from pathlib import Path
from typing import List

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def get_slide_images(images_folder: Path, pattern: str = "slide_*.png") -> List[Path]:
    """Get sorted list of slide images from folder."""
    images = sorted(images_folder.glob(pattern))
    return images


def create_carousel_presentation(
    slide_images: List[Path],
    output_path: Path,
    slides_per_page: int = 4,
    slide_width_inches: float = 16.0,
    slide_height_inches: float = 9.0,
    add_titles: bool = True,
    add_borders: bool = True
):
    """
    Create a carousel PowerPoint presentation from slide images.
    
    Args:
        slide_images: List of paths to slide image files
        output_path: Path where output .pptx will be saved
        slides_per_page: Number of slide thumbnails per carousel page
        slide_width_inches: Width of presentation slides (default 16" for 32:9)
        slide_height_inches: Height of presentation slides (default 9" for 32:9)
        add_titles: Whether to add title to each carousel page
        add_borders: Whether to add borders around slide thumbnails
    """
    print(f"Creating carousel presentation...")
    print(f"Total slides: {len(slide_images)}")
    print(f"Slides per page: {slides_per_page}")
    
    # Create new presentation
    prs = Presentation()
    
    # Set slide dimensions (32:9 aspect ratio)
    prs.slide_width = Inches(slide_width_inches)
    prs.slide_height = Inches(slide_height_inches)
    
    # Calculate carousel pages needed
    total_slides = len(slide_images)
    carousel_pages = (total_slides + slides_per_page - 1) // slides_per_page
    
    print(f"Creating {carousel_pages} carousel pages...")
    
    # Available space for thumbnails
    title_height = Inches(1.0) if add_titles else Inches(0.2)
    margin = Inches(0.5)
    
    available_width = Inches(slide_width_inches) - (2 * margin)
    available_height = Inches(slide_height_inches) - title_height - margin
    
    for page_num in range(carousel_pages):
        # Create blank slide
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        # Add background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(245, 245, 245)  # Light gray
        
        # Add title if requested
        if add_titles:
            start_idx = page_num * slides_per_page + 1
            end_idx = min(start_idx + slides_per_page - 1, total_slides)
            
            title_box = slide.shapes.add_textbox(
                margin, Inches(0.3),
                available_width, Inches(0.6)
            )
            title_frame = title_box.text_frame
            title_frame.text = f"Slides {start_idx} - {end_idx}"
            title_frame.paragraphs[0].font.size = Pt(32)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor(50, 50, 50)
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Calculate which slides go on this page
        page_start_idx = page_num * slides_per_page
        page_end_idx = min(page_start_idx + slides_per_page, total_slides)
        page_slide_images = slide_images[page_start_idx:page_end_idx]
        
        # Determine grid layout
        num_slides = len(page_slide_images)
        
        if num_slides <= 2:
            # Horizontal arrangement
            cols, rows = num_slides, 1
        elif num_slides <= 4:
            # 2x2 grid
            cols, rows = 2, 2
        elif num_slides <= 6:
            # 2x3 grid
            cols, rows = 3, 2
        else:
            # 3x3 grid for up to 9 slides
            cols, rows = 3, 3
        
        # Calculate thumbnail dimensions with spacing
        spacing = Inches(0.3)
        thumb_width = (available_width - spacing * (cols - 1)) / cols
        thumb_height = (available_height - spacing * (rows - 1)) / rows
        
        # Add thumbnail images
        for idx, img_path in enumerate(page_slide_images):
            row = idx // cols
            col = idx % cols
            
            left = margin + col * (thumb_width + spacing)
            top = title_height + row * (thumb_height + spacing)
            
            try:
                # Add image
                pic = slide.shapes.add_picture(
                    str(img_path),
                    left, top,
                    width=thumb_width,
                    height=thumb_height
                )
                
                # Add border if requested
                if add_borders:
                    line = pic.line
                    line.color.rgb = RGBColor(180, 180, 180)
                    line.width = Pt(2)
                
                # Add shadow (simplified - python-pptx has limited shadow support)
                try:
                    shadow = pic.shadow
                    shadow.inherit = False
                except:
                    pass  # Shadow formatting not fully supported
                
                # Add slide number label
                label_box = slide.shapes.add_textbox(
                    left, top + thumb_height - Inches(0.4),
                    Inches(0.6), Inches(0.35)
                )
                label_frame = label_box.text_frame
                label_frame.text = str(page_start_idx + idx + 1)
                label_frame.paragraphs[0].font.size = Pt(18)
                label_frame.paragraphs[0].font.bold = True
                label_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # Add background to label
                label_shape = label_box
                label_fill = label_shape.fill
                label_fill.solid()
                label_fill.fore_color.rgb = RGBColor(50, 50, 150)
                
            except Exception as e:
                print(f"  Warning: Could not add image {img_path.name}: {e}")
        
        print(f"  Created page {page_num + 1}/{carousel_pages} with {num_slides} slides")
    
    # Save presentation
    prs.save(str(output_path))
    print(f"\n✅ Carousel presentation created successfully!")
    print(f"   Output: {output_path}")
    print(f"   Total pages: {carousel_pages}")


def parse_args():
    parser = argparse.ArgumentParser(
        description="Build carousel PowerPoint from slide images"
    )
    parser.add_argument(
        "images_folder",
        help="Path to folder containing slide images"
    )
    parser.add_argument(
        "output",
        help="Output PowerPoint file path (.pptx)"
    )
    parser.add_argument(
        "-s", "--slides-per-page",
        type=int,
        default=4,
        help="Number of slide thumbnails per carousel page (default: 4)"
    )
    parser.add_argument(
        "--width",
        type=float,
        default=16.0,
        help="Slide width in inches (default: 16.0 for 32:9)"
    )
    parser.add_argument(
        "--height",
        type=float,
        default=9.0,
        help="Slide height in inches (default: 9.0 for 32:9)"
    )
    parser.add_argument(
        "--no-titles",
        action="store_true",
        help="Don't add titles to carousel pages"
    )
    parser.add_argument(
        "--no-borders",
        action="store_true",
        help="Don't add borders around thumbnails"
    )
    parser.add_argument(
        "--pattern",
        default="slide_*.png",
        help="File pattern for slide images (default: slide_*.png)"
    )
    return parser.parse_args()


def main():
    args = parse_args()
    
    images_folder = Path(args.images_folder)
    output_path = Path(args.output)
    
    # Validate inputs
    if not images_folder.exists():
        print(f"Error: Images folder not found: {images_folder}")
        return 1
    
    if not images_folder.is_dir():
        print(f"Error: Images path is not a folder: {images_folder}")
        return 1
    
    # Get slide images
    slide_images = get_slide_images(images_folder, args.pattern)
    
    if not slide_images:
        print(f"Error: No images found in {images_folder} matching pattern '{args.pattern}'")
        return 1
    
    print(f"Found {len(slide_images)} slide images")
    
    # Create carousel
    try:
        create_carousel_presentation(
            slide_images,
            output_path,
            slides_per_page=args.slides_per_page,
            slide_width_inches=args.width,
            slide_height_inches=args.height,
            add_titles=not args.no_titles,
            add_borders=not args.no_borders
        )
    except Exception as e:
        print(f"Error creating carousel: {e}")
        import traceback
        traceback.print_exc()
        return 1
    
    return 0


if __name__ == "__main__":
    exit(main())
