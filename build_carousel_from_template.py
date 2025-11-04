"""
Build a carousel PowerPoint presentation based on a template design.

This script creates a carousel-style presentation matching the design
from "Carousel Presentation Template.pptx" where slides scroll horizontally.

Usage:
    python build_carousel_from_template.py template.json images_folder/ output.pptx
"""

import argparse
import json
from pathlib import Path
from typing import Dict, Any, List

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


def load_template_definition(template_path: Path) -> Dict[str, Any]:
    """Load template definition from JSON file."""
    return json.loads(template_path.read_text(encoding='utf-8'))


def analyze_template_layout(template_def: Dict[str, Any]) -> Dict[str, Any]:
    """
    Analyze template to extract layout parameters.
    
    Returns:
        dict with slide_width, slide_height, picture_positions, decorative_shapes
    """
    metadata = template_def['metadata']
    
    # Get first slide to analyze layout
    first_slide = template_def['slides'][0]
    
    # Find all picture shapes and their positions
    pictures = [s for s in first_slide['shapes'] if s.get('is_picture', False)]
    rectangles = [s for s in first_slide['shapes'] if s['type'] == 'auto_shape']
    
    # Sort pictures by left position to understand the carousel order
    pictures_sorted = sorted(pictures, key=lambda p: p['left_inches'])
    
    layout_info = {
        'slide_width_inches': metadata['slide_width_inches'],
        'slide_height_inches': metadata['slide_height_inches'],
        'pictures': pictures_sorted,
        'rectangles': rectangles,
        'center_picture_index': 1 if len(pictures_sorted) >= 3 else 0,  # Middle one is center
    }
    
    return layout_info


def analyze_all_template_slides(template_def: Dict[str, Any]) -> List[List[Dict[str, Any]]]:
    """
    Analyze all template slides to extract the carousel pattern.
    
    Returns:
        List of lists, where each inner list contains picture definitions for that template slide
    """
    all_slides_pictures = []
    
    for slide_def in template_def['slides']:
        # Find all picture shapes and sort by left position
        pictures = [s for s in slide_def['shapes'] if s.get('is_picture', False)]
        pictures_sorted = sorted(pictures, key=lambda p: p['left_inches'])
        all_slides_pictures.append(pictures_sorted)
    
    return all_slides_pictures


def create_carousel_from_template(
    template_def: Dict[str, Any],
    slide_images: List[Path],
    output_path: Path
):
    """
    Create carousel presentation matching template design.
    
    The carousel progresses through slides, centering each one in sequence:
    - Slide 1 centers image 0
    - Slide 2 centers image 1
    - Slide 3 centers image 2
    - etc.
    
    Each page shows surrounding images in smaller sizes on left and right.
    
    Args:
        template_def: Template definition dict
        slide_images: List of slide image paths
        output_path: Output .pptx path
    """
    print(f"Creating carousel from template...")
    print(f"Template: {template_def['source_file']}")
    print(f"Total slides to include: {len(slide_images)}")
    
    # Get metadata
    metadata = template_def['metadata']
    slide_width = metadata['slide_width_inches']
    slide_height = metadata['slide_height_inches']
    
    # Analyze all template slides to get the pattern
    template_patterns = analyze_all_template_slides(template_def)
    
    print(f"Template layout: {slide_width:.2f}\" x {slide_height:.2f}\"")
    print(f"Template has {len(template_patterns)} example slides showing the pattern")
    
        # Get decorative rectangles from first template slide
    rectangles = [s for s in template_def['slides'][0]['shapes'] if s['type'] == 'auto_shape']
    
    # Sort rectangles by left position to identify left and right panels
    rectangles_sorted = sorted(rectangles, key=lambda r: r['left_inches'])    # Create new presentation
    prs = Presentation()
    prs.slide_width = Inches(slide_width)
    prs.slide_height = Inches(slide_height)
    
    # Each carousel page centers one slide
    # We need one page per input slide
    total_pages = len(slide_images)
    
    print(f"Creating {total_pages} carousel pages...")
    print(f"Pattern: Each page centers a different slide with surrounding slides visible")
    
    for center_idx in range(total_pages):
        # Create blank slide
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        # Set background to black
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)
        
        # Determine which template pattern to use
        # Use template slide pattern based on center position
        # Template slide 0: centers slide 0 (shows 0,1,2)
        # Template slide 1: centers slide 1 (shows 0,1,2,3)
        # Template slide 2: centers slide 2 (shows 0,1,2,3,4)
        # Template slide 3: centers slide 3 (shows 1,2,3,4,5)
        # Template slide 4: centers slide 4 (shows 2,3,4,5,6)
        
        # For our carousel, we use the appropriate template pattern
        # If we're centering slide N:
        # - Use template pattern based on position (early, middle, or late in sequence)
        
        # Determine which template slide pattern to use
        if center_idx == 0:
            # First slide: use template slide 0 pattern (3 images: center, right, far-right)
            template_idx = 0
            visible_range = range(0, min(3, total_pages))
        elif center_idx == 1 and len(template_patterns) > 1:
            # Second slide: use template slide 1 pattern (4 images: left, center, right, far-right)
            template_idx = 1
            visible_range = range(0, min(4, total_pages))
        elif center_idx >= 2 and center_idx < total_pages - 2:
            # Middle slides: use template slide 2+ pattern (5 images: far-left, left, center, right, far-right)
            template_idx = min(2, len(template_patterns) - 1)
            visible_range = range(max(0, center_idx - 2), min(center_idx + 3, total_pages))
        elif center_idx == total_pages - 2 and len(template_patterns) > 3:
            # Second to last: similar to slide 1 but reversed
            template_idx = min(3, len(template_patterns) - 1)
            visible_range = range(max(0, center_idx - 2), total_pages)
        else:
            # Last slide(s): use last template pattern or pattern 2
            template_idx = min(len(template_patterns) - 1, 2)
            visible_range = range(max(0, center_idx - 2), total_pages)
        
        # Get the picture positions from the selected template
        template_pictures = template_patterns[template_idx]
        
        # Map visible slides to template positions
        visible_slides = list(visible_range)
        
        # Find which position in visible_slides is our center
        center_position_in_visible = visible_slides.index(center_idx) if center_idx in visible_slides else 0
        
        # Add images according to template positions
        for pic_idx, pic_def in enumerate(template_pictures):
            if pic_idx >= len(visible_slides):
                break
            
            slide_img_idx = visible_slides[pic_idx]
            
            if slide_img_idx < 0 or slide_img_idx >= len(slide_images):
                continue
            
            img_path = slide_images[slide_img_idx]
            
            try:
                pic = slide.shapes.add_picture(
                    str(img_path),
                    Inches(pic_def['left_inches']),
                    Inches(pic_def['top_inches']),
                    width=Inches(pic_def['width_inches']),
                    height=Inches(pic_def['height_inches'])
                )
                
                # Add border
                line = pic.line
                line.color.rgb = RGBColor(200, 200, 200)
                line.width = Pt(1)
                
            except Exception as e:
                print(f"  Warning: Could not add image {img_path.name}: {e}")
        
        # Add gradient overlay rectangles AFTER images (so they appear on top)
        # Left panel: gradient from transparent (right edge) to black (left edge)
        # Right panel: gradient from transparent (left edge) to black (right edge)
        
        if len(rectangles_sorted) >= 2:
            # Left rectangle (index 0 - leftmost)
            left_rect_def = rectangles_sorted[0]
            try:
                left_rect = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(left_rect_def['left_inches']),
                    Inches(left_rect_def['top_inches']),
                    Inches(left_rect_def['width_inches']),
                    Inches(left_rect_def['height_inches'])
                )
                left_rect.rotation = left_rect_def.get('rotation', 0.0)
                
                # Gradient from black (left edge) to transparent (right/center)
                left_fill = left_rect.fill
                left_fill.gradient()
                left_fill.gradient_angle = 0  # Left to right gradient
                
                # Access gradient stops - need to set positions and alpha
                # Stop 0: Black at left edge (position 0.0)
                left_fill.gradient_stops[0].position = 0.0
                left_fill.gradient_stops[0].color.rgb = RGBColor(0, 0, 0)
                
                # Stop 1: Transparent at right edge (position 1.0)
                left_fill.gradient_stops[1].position = 1.0
                left_fill.gradient_stops[1].color.rgb = RGBColor(0, 0, 0)
                
                # Try to set alpha/transparency via the internal XML if available
                try:
                    # Access the gradient stop's XML element to set alpha directly
                    gs1_element = left_fill.gradient_stops[1]._element
                    # Find or create the alpha element in the color
                    from lxml import etree
                    alpha_elem = gs1_element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                    if alpha_elem is None:
                        color_elem = gs1_element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                        if color_elem is not None:
                            alpha_elem = etree.SubElement(color_elem, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                    if alpha_elem is not None:
                        alpha_elem.set('val', '0')  # 0 = fully transparent
                except Exception as inner_e:
                    print(f"    Could not set alpha for left gradient: {inner_e}")
                
                # Remove line
                left_rect.line.fill.background()
                
            except Exception as e:
                print(f"  Warning: Could not add left gradient rectangle: {e}")
            
            # Right rectangle (index 1 - rightmost)
            right_rect_def = rectangles_sorted[1]
            try:
                right_rect = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(right_rect_def['left_inches']),
                    Inches(right_rect_def['top_inches']),
                    Inches(right_rect_def['width_inches']),
                    Inches(right_rect_def['height_inches'])
                )
                right_rect.rotation = right_rect_def.get('rotation', 0.0)
                
                # Gradient from transparent (left/center) to black (right edge)
                right_fill = right_rect.fill
                right_fill.gradient()
                right_fill.gradient_angle = 180  # Right to left gradient
                
                # Stop 0: Transparent at left edge (position 0.0) 
                right_fill.gradient_stops[0].position = 0.0
                right_fill.gradient_stops[0].color.rgb = RGBColor(0, 0, 0)
                
                # Try to set alpha for transparent stop
                try:
                    from lxml import etree
                    gs0_element = right_fill.gradient_stops[0]._element
                    alpha_elem = gs0_element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                    if alpha_elem is None:
                        color_elem = gs0_element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                        if color_elem is not None:
                            alpha_elem = etree.SubElement(color_elem, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                    if alpha_elem is not None:
                        alpha_elem.set('val', '0')  # 0 = fully transparent
                except Exception as inner_e:
                    print(f"    Could not set alpha for right gradient: {inner_e}")
                
                # Stop 1: Black at right edge (position 1.0)
                right_fill.gradient_stops[1].position = 1.0
                right_fill.gradient_stops[1].color.rgb = RGBColor(0, 0, 0)
                
                # Remove line
                right_rect.line.fill.background()
                
            except Exception as e:
                print(f"  Warning: Could not add right gradient rectangle: {e}")
        

        
        print(f"  Created page {center_idx + 1}/{total_pages} (centering slide {center_idx + 1})")
    
    # Save presentation
    prs.save(str(output_path))
    print(f"\n✅ Carousel presentation created successfully!")
    print(f"   Output: {output_path}")
    print(f"   Total pages: {total_pages}")


def get_slide_images(images_folder: Path, pattern: str = "slide_*.png") -> List[Path]:
    """Get sorted list of slide images from folder."""
    return sorted(images_folder.glob(pattern))


def parse_args():
    parser = argparse.ArgumentParser(
        description="Build carousel PowerPoint from template and slide images"
    )
    parser.add_argument(
        "template",
        help="Path to template definition file (JSON)"
    )
    parser.add_argument(
        "images_folder",
        help="Path to folder containing slide images"
    )
    parser.add_argument(
        "output",
        help="Output PowerPoint file path (.pptx)"
    )
    parser.add_argument(
        "--pattern",
        default="slide_*.png",
        help="File pattern for slide images (default: slide_*.png)"
    )
    return parser.parse_args()


def main():
    args = parse_args()
    
    template_path = Path(args.template)
    images_folder = Path(args.images_folder)
    output_path = Path(args.output)
    
    # Validate inputs
    if not template_path.exists():
        print(f"Error: Template file not found: {template_path}")
        return 1
    
    if not images_folder.exists():
        print(f"Error: Images folder not found: {images_folder}")
        return 1
    
    # Load template
    try:
        template_def = load_template_definition(template_path)
    except Exception as e:
        print(f"Error loading template: {e}")
        return 1
    
    # Get slide images
    slide_images = get_slide_images(images_folder, args.pattern)
    
    if not slide_images:
        print(f"Error: No images found in {images_folder} matching pattern '{args.pattern}'")
        return 1
    
    print(f"Found {len(slide_images)} slide images")
    
    # Create carousel
    try:
        create_carousel_from_template(template_def, slide_images, output_path)
    except Exception as e:
        print(f"Error creating carousel: {e}")
        import traceback
        traceback.print_exc()
        return 1
    
    return 0


if __name__ == "__main__":
    exit(main())
