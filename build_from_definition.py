"""
Build a PowerPoint presentation from a YAML/JSON definition file and slide images.

This script creates a new presentation using:
1. A definition file (YAML or JSON) created by pptx_to_definition.py
2. Slide images exported from the original presentation

Usage:
    python build_from_definition.py definition.yaml images_folder/ output.pptx
"""

import argparse
import json
from pathlib import Path
from typing import Dict, Any

from pptx import Presentation
from pptx.util import Inches

try:
    import yaml
    YAML_AVAILABLE = True
except ImportError:
    YAML_AVAILABLE = False


def load_definition(definition_path: Path) -> Dict[str, Any]:
    """Load presentation definition from YAML or JSON file."""
    content = definition_path.read_text(encoding='utf-8')
    
    if definition_path.suffix.lower() in ['.yaml', '.yml']:
        if not YAML_AVAILABLE:
            raise RuntimeError("PyYAML not installed. Install with: pip install PyYAML")
        return yaml.safe_load(content)
    elif definition_path.suffix.lower() == '.json':
        return json.loads(content)
    else:
        raise ValueError(f"Unsupported file format: {definition_path.suffix}. Use .yaml, .yml, or .json")


def create_presentation_from_definition(
    definition: Dict[str, Any],
    images_folder: Path,
    output_path: Path,
    slide_width_inches: float = 16.0,
    slide_height_inches: float = 9.0
):
    """
    Create a new PowerPoint presentation from definition and slide images.
    
    Args:
        definition: Presentation definition dict
        images_folder: Path to folder containing slide images
        output_path: Path where output .pptx will be saved
        slide_width_inches: Width of slides (default 16" for 32:9)
        slide_height_inches: Height of slides (default 9" for 32:9)
    """
    print(f"Creating presentation from definition...")
    print(f"Source: {definition['source_file']}")
    print(f"Total slides in definition: {len(definition['slides'])}")
    print(f"Images folder: {images_folder}")
    
    # Create new presentation
    prs = Presentation()
    
    # Set slide dimensions (32:9 aspect ratio by default)
    prs.slide_width = Inches(slide_width_inches)
    prs.slide_height = Inches(slide_height_inches)
    
    # Process each slide
    slides_created = 0
    for slide_def in definition['slides']:
        slide_index = slide_def['index']
        slide_image_file = slide_def.get('slide_image')
        
        if not slide_image_file:
            print(f"  Warning: Slide {slide_index} has no image file specified, skipping...")
            continue
        
        image_path = images_folder / slide_image_file
        
        if not image_path.exists():
            print(f"  Warning: Image not found: {image_path}, skipping slide {slide_index}...")
            continue
        
        # Create blank slide
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)
        
        # Add slide image to fill the entire slide
        left = Inches(0)
        top = Inches(0)
        width = prs.slide_width
        height = prs.slide_height
        
        try:
            slide.shapes.add_picture(
                str(image_path),
                left, top,
                width=width,
                height=height
            )
            slides_created += 1
            print(f"  Added slide {slide_index}: {slide_image_file}")
        except Exception as e:
            print(f"  Error adding slide {slide_index}: {e}")
    
    # Save presentation
    prs.save(str(output_path))
    print(f"\n✅ Presentation created successfully!")
    print(f"   Output: {output_path}")
    print(f"   Slides created: {slides_created}")


def parse_args():
    parser = argparse.ArgumentParser(
        description="Build PowerPoint presentation from definition file and slide images"
    )
    parser.add_argument(
        "definition",
        help="Path to definition file (.yaml, .yml, or .json)"
    )
    parser.add_argument(
        "images_folder",
        help="Path to folder containing slide images"
    )
    parser.add_argument(
        "output",
        help="Output PowerPoint file path (.pptx)"
    )
    parser.add_argument(
        "--width",
        type=float,
        default=16.0,
        help="Slide width in inches (default: 16.0 for 32:9)"
    )
    parser.add_argument(
        "--height",
        type=float,
        default=9.0,
        help="Slide height in inches (default: 9.0 for 32:9)"
    )
    return parser.parse_args()


def main():
    args = parse_args()
    
    definition_path = Path(args.definition)
    images_folder = Path(args.images_folder)
    output_path = Path(args.output)
    
    # Validate inputs
    if not definition_path.exists():
        print(f"Error: Definition file not found: {definition_path}")
        return 1
    
    if not images_folder.exists():
        print(f"Error: Images folder not found: {images_folder}")
        return 1
    
    if not images_folder.is_dir():
        print(f"Error: Images path is not a folder: {images_folder}")
        return 1
    
    # Load definition
    try:
        definition = load_definition(definition_path)
    except Exception as e:
        print(f"Error loading definition: {e}")
        return 1
    
    # Create presentation
    try:
        create_presentation_from_definition(
            definition,
            images_folder,
            output_path,
            args.width,
            args.height
        )
    except Exception as e:
        print(f"Error creating presentation: {e}")
        import traceback
        traceback.print_exc()
        return 1
    
    return 0


if __name__ == "__main__":
    exit(main())
