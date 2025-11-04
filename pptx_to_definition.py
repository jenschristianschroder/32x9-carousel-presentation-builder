import argparse
import base64
import json
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu
from pptx.dml.color import RGBColor

try:
    import yaml  # Optional (PyYAML)
    YAML_AVAILABLE = True
except ImportError:  # pragma: no cover
    YAML_AVAILABLE = False

# -----------------------------
# Utility helpers
# -----------------------------

def parse_slide_range(range_str: str, total_slides: int) -> List[int]:
    """
    Parse slide range string and return list of slide indices (1-based).
    
    Examples:
        "3-12"      -> slides 3 to 12
        "3.."       -> slides 3 to last
        "..5"       -> slides 1 to 5
        "..5,7-9"   -> slides 1-5 and 7-9
        "1,3,5"     -> slides 1, 3, and 5
    """
    if not range_str:
        return list(range(1, total_slides + 1))
    
    slide_indices = set()
    
    # Split by comma for multiple ranges
    parts = [p.strip() for p in range_str.split(',')]
    
    for part in parts:
        if '..' in part:
            # Handle range with ".."
            start_str, end_str = part.split('..', 1)
            start = int(start_str) if start_str else 1
            end = int(end_str) if end_str else total_slides
            slide_indices.update(range(start, end + 1))
        elif '-' in part:
            # Handle range with "-"
            start_str, end_str = part.split('-', 1)
            start = int(start_str)
            end = int(end_str)
            slide_indices.update(range(start, end + 1))
        else:
            # Single slide number
            slide_indices.add(int(part))
    
    # Filter to valid range and sort
    valid_indices = [i for i in sorted(slide_indices) if 1 <= i <= total_slides]
    return valid_indices


def emu_to_points(value: Emu) -> float:
    return float(value) / 12700.0  # 1 point = 12700 EMUs


def emu_to_inches(value: Emu) -> float:
    return float(value) / 914400.0  # 1 inch = 914400 EMUs


def rgb_color_to_hex(color: Optional[RGBColor]) -> Optional[str]:
    if not color:
        return None
    try:
        # RGBColor is itself the tuple (r, g, b)
        r, g, b = color
        return f"#{r:02X}{g:02X}{b:02X}"
    except (TypeError, ValueError):
        return None


def safe_text(shape) -> str:
    if not hasattr(shape, "text"):
        return ""
    return shape.text or ""


def extract_paragraphs(shape) -> List[Dict[str, Any]]:
    paragraphs = []
    if not hasattr(shape, "text_frame") or shape.text_frame is None:
        return paragraphs
    for p in shape.text_frame.paragraphs:
        runs = []
        for r in p.runs:
            font = r.font
            runs.append(
                {
                    "text": r.text,
                    "font": {
                        "name": font.name,
                        "size_pt": getattr(font.size, "pt", None),
                        "bold": font.bold,
                        "italic": font.italic,
                        "underline": font.underline,
                        "color_hex": rgb_color_to_hex(getattr(font.color, "rgb", None)),
                    },
                }
            )
        paragraphs.append(
            {
                "alignment": str(p.alignment) if p.alignment else None,
                "runs": runs,
            }
        )
    return paragraphs


def export_slide_as_image(pptx_path: Path, slide_index: int, export_dir: Path) -> Optional[str]:
    """Export a single slide as PNG image using COM automation."""
    import win32com.client
    import time
    
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = True  # PowerPoint requires this to be True
        
        presentation = powerpoint.Presentations.Open(str(pptx_path.absolute()), ReadOnly=True, WithWindow=False)
        
        filename = f"slide_{slide_index:03d}.png"
        output_path = export_dir / filename
        
        # Export slide as PNG (2 = PNG format)
        presentation.Slides(slide_index).Export(str(output_path.absolute()), "PNG")
        
        presentation.Close()
        powerpoint.Quit()
        
        print(f"  Exported slide image: {filename}")
        return filename
        
    except Exception as e:
        print(f"  Warning: Failed to export slide {slide_index}: {e}")
        try:
            if 'presentation' in locals():
                presentation.Close()
            if 'powerpoint' in locals():
                powerpoint.Quit()
        except:
            pass
        return None


def shape_type_name(shape) -> str:
    t = shape.shape_type
    mapping = {
        MSO_SHAPE_TYPE.AUTO_SHAPE: "auto_shape",
        MSO_SHAPE_TYPE.CHART: "chart",
        MSO_SHAPE_TYPE.PICTURE: "picture",
        MSO_SHAPE_TYPE.PLACEHOLDER: "placeholder",
        MSO_SHAPE_TYPE.GROUP: "group",
        MSO_SHAPE_TYPE.LINE: "line",
        MSO_SHAPE_TYPE.TABLE: "table",
        MSO_SHAPE_TYPE.TEXT_BOX: "text_box",
        MSO_SHAPE_TYPE.MEDIA: "media",
        MSO_SHAPE_TYPE.FREEFORM: "freeform",
    }
    return mapping.get(t, str(t))


def extract_table(shape) -> Optional[Dict[str, Any]]:
    if not hasattr(shape, "table"):
        return None
    table = shape.table
    rows = []
    for r in range(len(table.rows)):
        row_cells = []
        for c in range(len(table.columns)):
            cell = table.cell(r, c)
            row_cells.append(
                {
                    "text": cell.text,
                    "row_span": getattr(cell, "row_span", None),
                    "col_span": getattr(cell, "col_span", None),
                }
            )
        rows.append(row_cells)
    return {
        "rows": rows,
        "row_count": len(table.rows),
        "column_count": len(table.columns),
    }


def extract_notes(slide) -> Optional[str]:
    if not hasattr(slide, "notes_slide"):
        return None
    try:
        notes_slide = slide.notes_slide
        if notes_slide and notes_slide.notes_text_frame:
            return notes_slide.notes_text_frame.text
    except Exception:  # pragma: no cover
        return None
    return None


def extract_theme_metadata(prs: Presentation) -> Dict[str, Any]:
    return {
        "slide_width_inches": round(emu_to_inches(prs.slide_width), 4),
        "slide_height_inches": round(emu_to_inches(prs.slide_height), 4),
        "slide_count": len(prs.slides),
    }

# -----------------------------
# Core extraction
# -----------------------------

def extract_presentation_definition(
    pptx_path: Path,
    export_images: bool = False,
    max_slides: Optional[int] = None,
    slide_range: Optional[str] = None,
) -> Dict[str, Any]:
    prs = Presentation(str(pptx_path))

    image_export_dir = pptx_path.parent / f"{pptx_path.stem}_images"
    if export_images:
        image_export_dir.mkdir(exist_ok=True)

    presentation_def: Dict[str, Any] = {
        "source_file": pptx_path.name,
        "metadata": extract_theme_metadata(prs),
        "slides": [],
    }

    # Determine which slides to process
    total_slides = len(prs.slides)
    if slide_range:
        slide_indices = parse_slide_range(slide_range, total_slides)
    elif max_slides:
        slide_indices = list(range(1, min(max_slides + 1, total_slides + 1)))
    else:
        slide_indices = list(range(1, total_slides + 1))
    
    print(f"Processing {len(slide_indices)} of {total_slides} slides: {slide_indices if len(slide_indices) <= 10 else f'{slide_indices[:5]}...{slide_indices[-5:]}'}")

    for idx in slide_indices:
        slide = prs.slides[idx - 1]  # Convert to 0-based index
        
        print(f"Processing slide {idx}/{total_slides}...")
        
        # Export entire slide as image if requested
        slide_image_file = None
        if export_images:
            slide_image_file = export_slide_as_image(pptx_path, idx, image_export_dir)
        
        slide_info: Dict[str, Any] = {
            "index": idx,
            "layout_name": getattr(slide.slide_layout, "name", None),
            "slide_image": slide_image_file,
            "shapes": [],
            "notes": extract_notes(slide),
        }

        for shape in slide.shapes:
            shape_def: Dict[str, Any] = {
                "id": getattr(shape, "shape_id", None),
                "name": shape.name,
                "type": shape_type_name(shape),
                "left_inches": round(emu_to_inches(shape.left), 4),
                "top_inches": round(emu_to_inches(shape.top), 4),
                "width_inches": round(emu_to_inches(shape.width), 4),
                "height_inches": round(emu_to_inches(shape.height), 4),
                "rotation": getattr(shape, "rotation", 0.0),
                "has_text_frame": hasattr(shape, "text_frame") and shape.text_frame is not None,
            }

            if hasattr(shape, "text_frame") and shape.text_frame:
                shape_def["text"] = safe_text(shape)
                shape_def["paragraphs"] = extract_paragraphs(shape)

            if hasattr(shape, "has_table") and shape.has_table:
                table_def = extract_table(shape)
                if table_def:
                    shape_def["table"] = table_def

            # Note: We export entire slides as images, not individual picture shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shape_def["is_picture"] = True

            if shape.shape_type == MSO_SHAPE_TYPE.GROUP and hasattr(shape, "shapes"):
                group_items = []
                for gshape in shape.shapes:
                    group_items.append(
                        {
                            "name": gshape.name,
                            "type": shape_type_name(gshape),
                            "left_inches": round(emu_to_inches(gshape.left), 4),
                            "top_inches": round(emu_to_inches(gshape.top), 4),
                            "width_inches": round(emu_to_inches(gshape.width), 4),
                            "height_inches": round(emu_to_inches(gshape.height), 4),
                            "text": safe_text(gshape) if hasattr(gshape, "text_frame") else None,
                        }
                    )
                shape_def["group_children"] = group_items

            if shape.shape_type == MSO_SHAPE_TYPE.CHART and hasattr(shape, "chart"):
                chart = shape.chart
                shape_def["chart"] = {
                    "chart_type": str(chart.chart_type),
                    "has_legend": chart.has_legend,
                    "series_count": len(chart.series),
                    "categories_count": len(chart.plots[0].categories)
                    if chart.plots and chart.plots[0].categories
                    else None,
                }

            slide_info["shapes"].append(shape_def)

        presentation_def["slides"].append(slide_info)

    return presentation_def

# -----------------------------
# Serialization
# -----------------------------

def serialize_definition(definition: Dict[str, Any], fmt: str) -> str:
    fmt = fmt.lower()
    if fmt == "json":
        return json.dumps(definition, indent=2, ensure_ascii=False)
    elif fmt == "yaml":
        if not YAML_AVAILABLE:
            raise RuntimeError("PyYAML not installed. Install with: pip install PyYAML")
        return yaml.safe_dump(definition, sort_keys=False, allow_unicode=True)
    else:
        raise ValueError("Unsupported format. Use 'json' or 'yaml'.")

# -----------------------------
# CLI
# -----------------------------

def parse_args():
    parser = argparse.ArgumentParser(
        description="Extract structural definition from a PowerPoint (.pptx) file into JSON or YAML."
    )
    parser.add_argument("input", help="Path to input .pptx file")
    parser.add_argument(
        "-o",
        "--output",
        help="Output file path (defaults to <input>_definition.<ext>)",
    )
    parser.add_argument(
        "-f",
        "--format",
        choices=["json", "yaml"],
        default="json",
        help="Output format (default: json)",
    )
    parser.add_argument(
        "--export-images",
        action="store_true",
        help="Export each slide as a PNG image to a folder",
    )
    parser.add_argument(
        "--max-slides", type=int, help="Limit number of slides processed (deprecated, use --range)"
    )
    parser.add_argument(
        "--range",
        help="Slide range to export (e.g., '3-12', '3..', '..5', '..5,7-9', '1,3,5')"
    )
    parser.add_argument(
        "--pretty", action="store_true", help="Pretty print (JSON only)"
    )
    return parser.parse_args()


def main():
    args = parse_args()
    pptx_path = Path(args.input)
    if not pptx_path.exists():
        raise FileNotFoundError(f"Input file not found: {pptx_path}")

    if args.output:
        out_path = Path(args.output)
    else:
        ext = args.format
        out_path = pptx_path.parent / f"{pptx_path.stem}_definition.{ext}"

    definition = extract_presentation_definition(
        pptx_path=pptx_path,
        export_images=args.export_images,
        max_slides=args.max_slides,
        slide_range=args.range,
    )

    text = serialize_definition(definition, args.format)
    if args.pretty and args.format == "json":
        # Already pretty due to indent=2
        pass

    out_path.write_text(text, encoding="utf-8")
    print(f"Definition written to: {out_path}")
    if args.export_images:
        print("Images exported to folder:", f"{pptx_path.stem}_images")


if __name__ == "__main__":
    main()
